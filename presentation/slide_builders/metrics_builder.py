"""
Metrics and analytics slide builder.

This module creates slides with data visualizations and key metrics.
"""

from typing import Dict, Any, List, Optional, Union
from dataclasses import dataclass
from .base_builder import BaseSlideBuilder


@dataclass
class Metric:
    """Represents a single metric item."""
    label: str
    value: Union[int, float, str]
    unit: Optional[str] = None
    change: Optional[float] = None  # Percentage change
    trend: Optional[str] = None  # 'up', 'down', 'stable'
    target: Optional[Union[int, float]] = None
    category: Optional[str] = None


@dataclass
class ChartData:
    """Represents chart data for visualization."""
    chart_type: str  # 'bar', 'line', 'pie', 'column'
    title: str
    data: Dict[str, Union[int, float]]
    x_axis_label: Optional[str] = None
    y_axis_label: Optional[str] = None


class MetricsSlideBuilder(BaseSlideBuilder):
    """Builder for creating metrics and data slides."""
    
    def supports_slide_type(self, slide_type: str) -> bool:
        """Check if this builder supports the given slide type."""
        return slide_type.lower() in ['metrics', 'data', 'analytics', 'stats', 'performance']
    
    def build_slide(self, slide_data: Dict[str, Any], presentation_engine: Any) -> bool:
        """
        Build a metrics slide.
        
        Args:
            slide_data: Dictionary containing slide information
            presentation_engine: Engine instance (COM or python-pptx)
            
        Returns:
            bool: True if slide was created successfully
        """
        try:
            # Validate slide data
            if not self.validate_slide_data(slide_data):
                return False
            
            # Extract metrics and chart data
            metrics = self._extract_metrics(slide_data)
            charts = self._extract_charts(slide_data)
            
            # Create slide using the presentation engine
            slide = presentation_engine.add_slide()
            
            # Set slide title
            title = self.get_slide_title(slide_data)
            if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title'):
                slide.shapes.title.text = title
            
            # Choose layout based on content
            if charts and metrics:
                self._create_mixed_layout(slide, metrics, charts)
            elif charts:
                self._create_chart_focused_layout(slide, charts)
            elif metrics:
                self._create_metrics_focused_layout(slide, metrics)
            else:
                self._create_default_layout(slide, slide_data)
            
            self.slides_created += 1
            return True
            
        except Exception as e:
            self.handle_builder_error(e, "build metrics slide")
            return False
    
    def _extract_metrics(self, slide_data: Dict[str, Any]) -> List[Metric]:
        """Extract metrics from slide data."""
        metrics = []
        
        if 'metrics' in slide_data:
            for item in slide_data['metrics']:
                if isinstance(item, dict):
                    metrics.append(Metric(
                        label=item.get('label', 'Metric'),
                        value=item.get('value', 0),
                        unit=item.get('unit'),
                        change=item.get('change'),
                        trend=item.get('trend'),
                        target=item.get('target'),
                        category=item.get('category')
                    ))
        
        elif 'data' in slide_data:
            data = slide_data['data']
            if isinstance(data, dict):
                for key, value in data.items():
                    metrics.append(Metric(
                        label=key,
                        value=value,
                        category='general'
                    ))
        
        return metrics
    
    def _extract_charts(self, slide_data: Dict[str, Any]) -> List[ChartData]:
        """Extract chart data from slide data."""
        charts = []
        
        if 'charts' in slide_data:
            for chart_info in slide_data['charts']:
                if isinstance(chart_info, dict):
                    charts.append(ChartData(
                        chart_type=chart_info.get('type', 'bar'),
                        title=chart_info.get('title', 'Chart'),
                        data=chart_info.get('data', {}),
                        x_axis_label=chart_info.get('x_label'),
                        y_axis_label=chart_info.get('y_label')
                    ))
        
        return charts
    
    def _create_mixed_layout(self, slide: object, metrics: List[Metric], 
                           charts: List[ChartData]) -> None:
        """Create layout with both metrics and charts."""
        if not hasattr(slide, 'shapes'):
            return
        
        # Clear default placeholders
        self._clear_default_placeholders(slide)
        
        # Top section: Key metrics cards
        self._create_metrics_cards(slide, metrics[:4], 1, 1.5, 8, 2)
        
        # Bottom section: Chart
        if charts:
            self._create_chart(slide, charts[0], 1, 4, 8, 3)
    
    def _create_chart_focused_layout(self, slide: object, charts: List[ChartData]) -> None:
        """Create layout focused on charts."""
        if not hasattr(slide, 'shapes'):
            return
        
        # Clear default placeholders
        self._clear_default_placeholders(slide)
        
        # Single large chart or multiple smaller charts
        if len(charts) == 1:
            self._create_chart(slide, charts[0], 1, 1.5, 8, 5)
        else:
            # Multiple charts in grid
            self._create_chart_grid(slide, charts)
    
    def _create_metrics_focused_layout(self, slide: object, metrics: List[Metric]) -> None:
        """Create layout focused on metrics."""
        if not hasattr(slide, 'shapes'):
            return
        
        # Clear default placeholders
        self._clear_default_placeholders(slide)
        
        if len(metrics) <= 4:
            # Single row of metric cards
            self._create_metrics_cards(slide, metrics, 1, 2.5, 8, 2)
        else:
            # Grid layout for many metrics
            self._create_metrics_grid(slide, metrics)
    
    def _create_metrics_cards(self, slide: object, metrics: List[Metric], 
                            x: float, y: float, width: float, height: float) -> None:
        """Create metric cards layout."""
        try:
            shapes = getattr(slide, 'shapes', None)
            if not shapes or not hasattr(shapes, 'add_shape'):
                return
            
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.util import Inches, Pt
            
            card_width = (width - 0.6) / len(metrics)  # With margins
            
            for i, metric in enumerate(metrics):
                card_x = x + (i * (card_width + 0.2))
                
                # Create card background
                card = shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(card_x), Inches(y),
                    Inches(card_width), Inches(height)
                )
                
                # Style the card
                if hasattr(card, 'fill'):
                    card.fill.solid()
                    card.fill.fore_color.rgb = self._get_theme_color('accent1')
                
                if hasattr(card, 'line'):
                    card.line.color.rgb = self._get_theme_color('accent2')
                    card.line.width = Inches(0.01)
                
                # Add content to card
                if hasattr(card, 'text_frame'):
                    text_frame = card.text_frame
                    text_frame.clear()
                    text_frame.margin_left = Inches(0.1)
                    text_frame.margin_right = Inches(0.1)
                    text_frame.margin_top = Inches(0.1)
                    text_frame.margin_bottom = Inches(0.1)
                    
                    # Value (large)
                    p = text_frame.paragraphs[0]
                    value_text = str(metric.value)
                    if metric.unit:
                        value_text += f" {metric.unit}"
                    p.text = value_text
                    p.font.bold = True
                    p.font.size = Pt(24)
                    p.font.color.rgb = self._get_theme_color('text1')
                    p.alignment = 1  # Center alignment
                    
                    # Label (smaller)
                    p = text_frame.add_paragraph()
                    p.text = metric.label
                    p.font.size = Pt(12)
                    p.font.color.rgb = self._get_theme_color('text2')
                    p.alignment = 1  # Center alignment
                    p.space_before = Pt(6)
                    
                    # Change indicator if available
                    if metric.change is not None:
                        p = text_frame.add_paragraph()
                        change_text = f"{metric.change:+.1f}%"
                        if metric.trend == 'up':
                            change_text = f"↑ {change_text}"
                        elif metric.trend == 'down':
                            change_text = f"↓ {change_text}"
                        
                        p.text = change_text
                        p.font.size = Pt(10)
                        p.font.color.rgb = self._get_trend_color(metric.trend)
                        p.alignment = 1  # Center alignment
                        p.space_before = Pt(4)
                        
        except Exception as e:
            # Fallback to simple text
            self._add_simple_metrics_text(slide, metrics, x, y)
    
    def _create_chart(self, slide: object, chart_data: ChartData, 
                     x: float, y: float, width: float, height: float) -> None:
        """Create a chart on the slide."""
        try:
            shapes = getattr(slide, 'shapes', None)
            if not shapes or not hasattr(shapes, 'add_chart'):
                # Fallback to table representation
                self._create_data_table(slide, chart_data, x, y, width, height)
                return
            
            from pptx.enum.chart import XL_CHART_TYPE
            from pptx.chart.data import CategoryChartData
            from pptx.util import Inches
            
            # Prepare chart data
            chart_data_obj = CategoryChartData()
            chart_data_obj.categories = list(chart_data.data.keys())
            chart_data_obj.add_series('Values', list(chart_data.data.values()))
            
            # Determine chart type
            chart_type_map = {
                'bar': XL_CHART_TYPE.BAR_CLUSTERED,
                'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
                'line': XL_CHART_TYPE.LINE,
                'pie': XL_CHART_TYPE.PIE,
            }
            
            chart_type = chart_type_map.get(chart_data.chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
            
            # Add chart to slide
            chart = shapes.add_chart(
                chart_type,
                Inches(x), Inches(y),
                Inches(width), Inches(height),
                chart_data_obj
            )
            
            # Customize chart
            if hasattr(chart, 'chart'):
                chart_obj = chart.chart
                
                # Set title
                if hasattr(chart_obj, 'chart_title'):
                    chart_obj.chart_title.text_frame.text = chart_data.title
                
                # Set axis labels
                if hasattr(chart_obj, 'category_axis') and chart_data.x_axis_label:
                    chart_obj.category_axis.axis_title.text_frame.text = chart_data.x_axis_label
                
                if hasattr(chart_obj, 'value_axis') and chart_data.y_axis_label:
                    chart_obj.value_axis.axis_title.text_frame.text = chart_data.y_axis_label
                
        except Exception as e:
            # Fallback to data table
            self._create_data_table(slide, chart_data, x, y, width, height)
    
    def _create_data_table(self, slide: object, chart_data: ChartData, 
                          x: float, y: float, width: float, height: float) -> None:
        """Create a data table as fallback for charts."""
        try:
            shapes = getattr(slide, 'shapes', None)
            if not shapes or not hasattr(shapes, 'add_table'):
                return
            
            from pptx.util import Inches
            
            # Create table
            rows = len(chart_data.data) + 1  # +1 for header
            cols = 2  # Label and Value columns
            
            table = shapes.add_table(
                rows, cols,
                Inches(x), Inches(y),
                Inches(width), Inches(height)
            )
            
            # Header row
            table.table.cell(0, 0).text = "Category"
            table.table.cell(0, 1).text = "Value"
            
            # Data rows
            for i, (label, value) in enumerate(chart_data.data.items()):
                table.table.cell(i + 1, 0).text = str(label)
                table.table.cell(i + 1, 1).text = str(value)
                
        except Exception:
            # Final fallback to text
            self._add_simple_chart_text(slide, chart_data, x, y)
    
    def _create_chart_grid(self, slide: object, charts: List[ChartData]) -> None:
        """Create a grid of charts."""
        if len(charts) <= 2:
            # Side by side
            for i, chart in enumerate(charts[:2]):
                self._create_chart(slide, chart, 1 + (i * 4), 2, 3.5, 4)
        else:
            # 2x2 grid
            positions = [(1, 1.5), (5, 1.5), (1, 4.5), (5, 4.5)]
            for i, chart in enumerate(charts[:4]):
                x, y = positions[i]
                self._create_chart(slide, chart, x, y, 3.5, 2.5)
    
    def _create_metrics_grid(self, slide: object, metrics: List[Metric]) -> None:
        """Create a grid layout for many metrics."""
        try:
            shapes = getattr(slide, 'shapes', None)
            if not shapes:
                return
            
            from pptx.util import Inches, Pt
            
            # Calculate grid dimensions
            cols = 3
            rows = (len(metrics) + cols - 1) // cols
            
            card_width = 2.5
            card_height = 1.5
            margin_x = 0.5
            margin_y = 0.3
            start_x = 1
            start_y = 1.5
            
            for i, metric in enumerate(metrics):
                row = i // cols
                col = i % cols
                
                x = start_x + (col * (card_width + margin_x))
                y = start_y + (row * (card_height + margin_y))
                
                # Create simple metric card
                if hasattr(shapes, 'add_textbox'):
                    textbox = shapes.add_textbox(
                        Inches(x), Inches(y),
                        Inches(card_width), Inches(card_height)
                    )
                    
                    text_frame = textbox.text_frame
                    text_frame.clear()
                    
                    # Value
                    p = text_frame.paragraphs[0]
                    p.text = str(metric.value)
                    if metric.unit:
                        p.text += f" {metric.unit}"
                    p.font.bold = True
                    p.font.size = Pt(18)
                    p.alignment = 1  # Center
                    
                    # Label
                    p = text_frame.add_paragraph()
                    p.text = metric.label
                    p.font.size = Pt(11)
                    p.alignment = 1  # Center
                    
        except Exception:
            pass
    
    def _create_default_layout(self, slide: object, slide_data: Dict[str, Any]) -> None:
        """Create default layout when no specific data is available."""
        try:
            content = slide_data.get('content', slide_data.get('points', []))
            if content:
                self._add_bullet_points(slide, content)
        except Exception:
            pass
    
    def _clear_default_placeholders(self, slide: object) -> None:
        """Remove default content placeholders."""
        try:
            if hasattr(slide, 'placeholders'):
                placeholders_to_remove = []
                for placeholder in getattr(slide, 'placeholders', []):
                    if hasattr(placeholder, 'placeholder_format'):
                        if placeholder.placeholder_format.type in [2, 7]:  # Content placeholders
                            placeholders_to_remove.append(placeholder)
                
                for placeholder in placeholders_to_remove:
                    try:
                        placeholder._element.getparent().remove(placeholder._element)
                    except Exception:
                        pass
        except Exception:
            pass
    
    def _add_bullet_points(self, slide: object, content: List[str]) -> None:
        """Add bullet points to slide."""
        try:
            if not content:
                return
            
            shapes = getattr(slide, 'shapes', None)
            if shapes and hasattr(shapes, 'add_textbox'):
                from pptx.util import Inches
                
                textbox = shapes.add_textbox(
                    Inches(1), Inches(2),
                    Inches(8), Inches(4)
                )
                
                text_frame = textbox.text_frame
                text_frame.clear()
                
                for i, point in enumerate(content):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = f"• {point}"
                    p.level = 0
        except Exception:
            pass
    
    def _get_theme_color(self, color_name: str) -> object:
        """Get theme color."""
        try:
            from pptx.dml.color import RGBColor
            
            colors = {
                'accent1': RGBColor(68, 114, 196),
                'accent2': RGBColor(112, 173, 71),
                'text1': RGBColor(68, 68, 68),
                'text2': RGBColor(89, 89, 89),
            }
            return colors.get(color_name, RGBColor(68, 68, 68))
        except ImportError:
            return None
    
    def _get_trend_color(self, trend: Optional[str]) -> object:
        """Get color for trend indicators."""
        try:
            from pptx.dml.color import RGBColor
            
            colors = {
                'up': RGBColor(34, 139, 34),     # Green
                'down': RGBColor(220, 20, 60),    # Red
                'stable': RGBColor(255, 165, 0),  # Orange
            }
            return colors.get(trend or 'stable', RGBColor(89, 89, 89))
        except ImportError:
            return None
    
    def _add_simple_metrics_text(self, slide: object, metrics: List[Metric], 
                               x: float, y: float) -> None:
        """Add simple text representation of metrics."""
        try:
            shapes = getattr(slide, 'shapes', None)
            if shapes and hasattr(shapes, 'add_textbox'):
                from pptx.util import Inches
                
                text_content = []
                for metric in metrics:
                    value_text = f"{metric.label}: {metric.value}"
                    if metric.unit:
                        value_text += f" {metric.unit}"
                    text_content.append(value_text)
                
                textbox = shapes.add_textbox(
                    Inches(x), Inches(y),
                    Inches(8), Inches(4)
                )
                textbox.text = "\n".join(text_content)
        except Exception:
            pass
    
    def _add_simple_chart_text(self, slide: object, chart_data: ChartData, 
                             x: float, y: float) -> None:
        """Add simple text representation of chart data."""
        try:
            shapes = getattr(slide, 'shapes', None)
            if shapes and hasattr(shapes, 'add_textbox'):
                from pptx.util import Inches
                
                text_content = [f"{chart_data.title}:"]
                for label, value in chart_data.data.items():
                    text_content.append(f"  {label}: {value}")
                
                textbox = shapes.add_textbox(
                    Inches(x), Inches(y),
                    Inches(6), Inches(3)
                )
                textbox.text = "\n".join(text_content)
        except Exception:
            pass
