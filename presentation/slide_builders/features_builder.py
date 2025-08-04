"""
Features slide builder.

This module creates slides highlighting key features and capabilities.
"""

from typing import Dict, Any, List, Optional
from dataclasses import dataclass
from .base_builder import BaseSlideBuilder


@dataclass
class Feature:
    """Represents a single feature item."""
    title: str
    description: str
    icon: Optional[str] = None
    category: Optional[str] = None
    priority: int = 1


class FeaturesSlideBuilder(BaseSlideBuilder):
    """Builder for creating feature highlight slides."""
    
    def supports_slide_type(self, slide_type: str) -> bool:
        """Check if this builder supports the given slide type."""
        return slide_type.lower() in ['features', 'highlights', 'capabilities', 'benefits']
    
    def build_slide(self, slide_data: Dict[str, Any], presentation_engine: Any) -> bool:
        """
        Build a features slide.
        
        Args:
            slide_data: Dictionary containing slide information
            presentation_engine: Engine instance (COM or python-pptx)
            
        Returns:
            bool: True if slide was created successfully
        """
        try:
            # Validate slide data
            if not slide_data:
                return False
            
            # Extract features data
            features = self._extract_features(slide_data)
            
            # Create slide using the presentation engine
            slide = presentation_engine.add_slide()
            
            # Set slide title
            title = slide_data.get('title', 'Key Features')
            if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title'):
                slide.shapes.title.text = title
            
            # Create layout based on number of features
            if len(features) <= 3:
                self._create_three_column_layout(slide, features)
            elif len(features) <= 6:
                self._create_grid_layout(slide, features)
            else:
                self._create_categorized_layout(slide, features)
            
            self.slides_created += 1
            return True
            
        except Exception as e:
            self.handle_builder_error(e, "build features slide")
            return False
    
    def _extract_features(self, slide_data: Dict[str, Any]) -> List[Feature]:
        """Extract and structure features from slide data."""
        features = []
        
        # Handle different data formats
        if 'features' in slide_data:
            for item in slide_data['features']:
                if isinstance(item, dict):
                    features.append(Feature(
                        title=item.get('title', 'Feature'),
                        description=item.get('description', ''),
                        icon=item.get('icon'),
                        category=item.get('category'),
                        priority=item.get('priority', 1)
                    ))
                else:
                    features.append(Feature(
                        title=str(item),
                        description='',
                        priority=1
                    ))
        
        elif 'points' in slide_data:
            for i, point in enumerate(slide_data['points']):
                features.append(Feature(
                    title=f"Feature {i+1}",
                    description=point,
                    priority=1
                ))
        
        elif 'content' in slide_data:
            content = slide_data['content']
            if isinstance(content, list):
                for i, item in enumerate(content):
                    features.append(Feature(
                        title=f"Feature {i+1}",
                        description=str(item),
                        priority=1
                    ))
        
        # Sort by priority
        features.sort(key=lambda x: x.priority, reverse=True)
        return features
    
    def _create_three_column_layout(self, slide: object, features: List[Feature]) -> None:
        """Create a three-column layout for features."""
        if not hasattr(slide, 'shapes') or not hasattr(slide, 'placeholders'):
            return
        
        # Remove default content placeholder if exists
        self._clear_default_placeholders(slide)
        
        # Calculate column positions
        slide_width = 10  # Standard slide width in inches
        slide_height = 7.5  # Standard slide height in inches
        
        col_width = (slide_width - 2) / 3  # 3 columns with margins
        col_height = 5
        start_x = 1
        start_y = 1.5
        
        for i, feature in enumerate(features[:3]):
            x = start_x + (i * (col_width + 0.2))
            
            # Create feature box
            self._create_feature_box(slide, feature, x, start_y, col_width, col_height)
    
    def _create_grid_layout(self, slide: object, features: List[Feature]) -> None:
        """Create a 2x3 grid layout for features."""
        if not hasattr(slide, 'shapes') or not hasattr(slide, 'placeholders'):
            return
        
        # Remove default content placeholder if exists
        self._clear_default_placeholders(slide)
        
        # Grid dimensions
        cols = 3
        rows = 2
        box_width = 3
        box_height = 2.5
        margin_x = 0.5
        margin_y = 0.3
        start_x = 1
        start_y = 1.5
        
        for i, feature in enumerate(features[:6]):
            row = i // cols
            col = i % cols
            
            x = start_x + (col * (box_width + margin_x))
            y = start_y + (row * (box_height + margin_y))
            
            self._create_feature_box(slide, feature, x, y, box_width, box_height)
    
    def _create_categorized_layout(self, slide: object, features: List[Feature]) -> None:
        """Create a categorized layout for many features."""
        if not hasattr(slide, 'shapes') or not hasattr(slide, 'placeholders'):
            return
        
        # Remove default content placeholder if exists
        self._clear_default_placeholders(slide)
        
        # Group features by category
        categories = {}
        for feature in features:
            category = feature.category or 'General'
            if category not in categories:
                categories[category] = []
            categories[category].append(feature)
        
        # Create sections for each category
        y_offset = 1.5
        for category, cat_features in categories.items():
            y_offset = self._create_category_section(slide, category, cat_features, y_offset)
    
    def _create_feature_box(self, slide: object, feature: Feature, x: float, y: float, 
                           width: float, height: float) -> None:
        """Create a single feature box."""
        try:
            # Create background shape
            shapes = getattr(slide, 'shapes', None)
            if shapes and hasattr(shapes, 'add_shape'):
                from pptx.enum.shapes import MSO_SHAPE
                from pptx.util import Inches
                
                # Add rounded rectangle
                shape = shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(x), Inches(y),
                    Inches(width), Inches(height)
                )
                
                # Style the shape
                if hasattr(shape, 'fill'):
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = self._get_theme_color('accent1')
                
                if hasattr(shape, 'line'):
                    shape.line.color.rgb = self._get_theme_color('accent2')
                    shape.line.width = Inches(0.02)
                
                # Add title text
                if hasattr(shape, 'text_frame'):
                    text_frame = shape.text_frame
                    text_frame.clear()
                    
                    # Feature title
                    p = text_frame.paragraphs[0]
                    p.text = feature.title
                    p.font.bold = True
                    p.font.size = 14  # Use int instead of Pt
                    p.font.color.rgb = self._get_theme_color('text1')
                    
                    # Feature description
                    if feature.description:
                        p = text_frame.add_paragraph()
                        p.text = feature.description
                        p.font.size = 11  # Use int instead of Pt
                        p.font.color.rgb = self._get_theme_color('text2')
                        # p.space_before = 6  # Skip space_before for compatibility
                
        except Exception as e:
            # Fallback to simple text if shapes not available
            self._add_simple_text(slide, f"{feature.title}: {feature.description}", x, y)
    
    def _create_category_section(self, slide: object, category: str, 
                               features: List[Feature], y_offset: float) -> float:
        """Create a section for a feature category."""
        try:
            shapes = getattr(slide, 'shapes', None)
            if shapes and hasattr(shapes, 'add_textbox'):
                from pptx.util import Inches, Pt
                
                # Category header
                header_box = shapes.add_textbox(
                    Inches(1), Inches(y_offset),
                    Inches(8), Inches(0.5)
                )
                header_frame = header_box.text_frame
                header_p = header_frame.paragraphs[0]
                header_p.text = category
                header_p.font.bold = True
                header_p.font.size = Pt(16)
                header_p.font.color.rgb = self._get_theme_color('accent1')
                
                y_offset += 0.7
                
                # Features in this category
                for feature in features[:4]:  # Limit to 4 per category
                    feature_box = shapes.add_textbox(
                        Inches(1.5), Inches(y_offset),
                        Inches(7), Inches(0.4)
                    )
                    feature_frame = feature_box.text_frame
                    feature_p = feature_frame.paragraphs[0]
                    feature_p.text = f"• {feature.title}"
                    if feature.description:
                        feature_p.text += f": {feature.description}"
                    feature_p.font.size = Pt(12)
                    feature_p.font.color.rgb = self._get_theme_color('text1')
                    
                    y_offset += 0.5
                
                y_offset += 0.3  # Space between categories
                
        except Exception:
            # Fallback
            y_offset += 1.5
        
        return y_offset
    
    def _clear_default_placeholders(self, slide: object) -> None:
        """Remove default content placeholders."""
        try:
            if hasattr(slide, 'placeholders'):
                for placeholder in getattr(slide, 'placeholders', []):
                    if hasattr(placeholder, 'placeholder_format'):
                        if placeholder.placeholder_format.type in [2, 7]:  # Content placeholders
                            try:
                                placeholder._element.getparent().remove(placeholder._element)
                            except Exception:
                                pass
        except Exception:
            pass
    
    def _get_theme_color(self, color_name: str) -> object:
        """Get theme color (fallback implementation)."""
        try:
            from pptx.dml.color import RGBColor
            
            colors = {
                'accent1': RGBColor(68, 114, 196),
                'accent2': RGBColor(112, 173, 71),
                'text1': RGBColor(68, 68, 68),
                'text2': RGBColor(89, 89, 89),
            }
            return colors.get(color_name, RGBColor(68, 68, 68))
        except ImportError:
            return None
    
    def _add_simple_text(self, slide: object, text: str, x: float, y: float) -> None:
        """Add simple text as fallback."""
        try:
            shapes = getattr(slide, 'shapes', None)
            if shapes and hasattr(shapes, 'add_textbox'):
                from pptx.util import Inches
                textbox = shapes.add_textbox(
                    Inches(x), Inches(y), Inches(3), Inches(1)
                )
                textbox.text = text
        except Exception:
            pass
