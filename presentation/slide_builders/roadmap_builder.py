"""
Roadmap slide builder.

This module creates slides with project timelines and roadmaps.
"""

from typing import Dict, Any, List, Optional
from dataclasses import dataclass
from datetime import datetime, date
from .base_builder import BaseSlideBuilder


@dataclass
class RoadmapItem:
    """Represents a single roadmap item."""
    title: str
    description: str
    start_date: Optional[date] = None
    end_date: Optional[date] = None
    status: str = 'planned'  # 'completed', 'in_progress', 'planned', 'delayed'
    priority: str = 'medium'  # 'high', 'medium', 'low'
    category: Optional[str] = None
    progress: int = 0  # 0-100 percentage


@dataclass
class Milestone:
    """Represents a milestone in the roadmap."""
    title: str
    date: date
    description: Optional[str] = None
    status: str = 'planned'


class RoadmapSlideBuilder(BaseSlideBuilder):
    """Builder for creating roadmap and timeline slides."""
    
    def supports_slide_type(self, slide_type: str) -> bool:
        """Check if this builder supports the given slide type."""
        return slide_type.lower() in ['roadmap', 'timeline', 'milestones', 'schedule', 'plan']
    
    def build_slide(self, slide_data: Dict[str, Any], presentation_engine: Any) -> bool:
        """
        Build a roadmap slide.
        
        Args:
            slide_data: Dictionary containing slide information
            presentation_engine: Engine instance (COM or python-pptx)
            
        Returns:
            bool: True if slide was created successfully
        """
        try:
            # Validate slide data
            if not slide_data:
                return False
            
            # Extract roadmap data
            roadmap_items = self._extract_roadmap_items(slide_data)
            milestones = self._extract_milestones(slide_data)
            
            # Create slide using the presentation engine
            slide = presentation_engine.add_slide()
            
            # Set slide title
            title = slide_data.get('title', 'Project Roadmap')
            if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title'):
                slide.shapes.title.text = title
            
            # Choose layout based on content
            if roadmap_items and milestones:
                self._create_combined_layout(slide, roadmap_items, milestones)
            elif roadmap_items:
                self._create_timeline_layout(slide, roadmap_items)
            elif milestones:
                self._create_milestones_layout(slide, milestones)
            else:
                self._create_default_layout(slide, slide_data)
            
            self.slides_created += 1
            return True
            
        except Exception as e:
            self.handle_builder_error(e, "build roadmap slide")
            return False
    
    def _extract_roadmap_items(self, slide_data: Dict[str, Any]) -> List[RoadmapItem]:
        """Extract roadmap items from slide data."""
        items = []
        
        if 'roadmap' in slide_data:
            for item_data in slide_data['roadmap']:
                if isinstance(item_data, dict):
                    items.append(RoadmapItem(
                        title=item_data.get('title', 'Task'),
                        description=item_data.get('description', ''),
                        start_date=self._parse_date(item_data.get('start_date')),
                        end_date=self._parse_date(item_data.get('end_date')),
                        status=item_data.get('status', 'planned'),
                        priority=item_data.get('priority', 'medium'),
                        category=item_data.get('category'),
                        progress=item_data.get('progress', 0)
                    ))
        
        elif 'timeline' in slide_data:
            for item_data in slide_data['timeline']:
                if isinstance(item_data, dict):
                    items.append(RoadmapItem(
                        title=item_data.get('title', 'Task'),
                        description=item_data.get('description', ''),
                        start_date=self._parse_date(item_data.get('date')),
                        status=item_data.get('status', 'planned'),
                        category=item_data.get('phase')
                    ))
        
        elif 'points' in slide_data:
            for i, point in enumerate(slide_data['points']):
                items.append(RoadmapItem(
                    title=f"Phase {i+1}",
                    description=str(point),
                    status='planned'
                ))
        
        return items
    
    def _extract_milestones(self, slide_data: Dict[str, Any]) -> List[Milestone]:
        """Extract milestones from slide data."""
        milestones = []
        
        if 'milestones' in slide_data:
            for milestone_data in slide_data['milestones']:
                if isinstance(milestone_data, dict):
                    milestones.append(Milestone(
                        title=milestone_data.get('title', 'Milestone'),
                        date=self._parse_date(milestone_data.get('date')) or date.today(),
                        description=milestone_data.get('description'),
                        status=milestone_data.get('status', 'planned')
                    ))
        
        return milestones
    
    def _parse_date(self, date_str: Any) -> Optional[date]:
        """Parse date string into date object."""
        if not date_str:
            return None
        
        if isinstance(date_str, date):
            return date_str
        
        if isinstance(date_str, datetime):
            return date_str.date()
        
        try:
            # Try common date formats
            for fmt in ['%Y-%m-%d', '%m/%d/%Y', '%d/%m/%Y', '%Y-%m', '%m/%Y']:
                try:
                    return datetime.strptime(str(date_str), fmt).date()
                except ValueError:
                    continue
        except Exception:
            pass
        
        return None
    
    def _create_combined_layout(self, slide: object, roadmap_items: List[RoadmapItem], 
                              milestones: List[Milestone]) -> None:
        """Create layout with both roadmap and milestones."""
        if not hasattr(slide, 'shapes'):
            return
        
        # Top section: Timeline
        self._create_timeline_section(slide, roadmap_items, 1, 1.5, 8, 3)
        
        # Bottom section: Milestones
        self._create_milestones_section(slide, milestones, 1, 5, 8, 2)
    
    def _create_timeline_layout(self, slide: object, roadmap_items: List[RoadmapItem]) -> None:
        """Create timeline-focused layout."""
        if not hasattr(slide, 'shapes'):
            return
        
        # Determine layout based on number of items
        if len(roadmap_items) <= 5:
            self._create_horizontal_timeline(slide, roadmap_items)
        else:
            self._create_vertical_timeline(slide, roadmap_items)
    
    def _create_milestones_layout(self, slide: object, milestones: List[Milestone]) -> None:
        """Create milestones-focused layout."""
        if not hasattr(slide, 'shapes'):
            return
        
        self._create_milestones_section(slide, milestones, 1, 1.5, 8, 5)
    
    def _create_horizontal_timeline(self, slide: object, roadmap_items: List[RoadmapItem]) -> None:
        """Create horizontal timeline layout."""
        try:
            shapes = getattr(slide, 'shapes', None)
            if not shapes or not hasattr(shapes, 'add_shape'):
                return
            
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.util import Inches, Pt
            
            # Timeline parameters
            timeline_y = 3
            timeline_width = 8
            timeline_start_x = 1
            item_width = timeline_width / len(roadmap_items)
            
            # Draw timeline base line
            line = shapes.add_connector(
                1,  # Straight line
                Inches(timeline_start_x), Inches(timeline_y),
                Inches(timeline_start_x + timeline_width), Inches(timeline_y)
            )
            
            if hasattr(line, 'line'):
                line.line.color.rgb = self._get_theme_color('accent1')
                line.line.width = Inches(0.05)
            
            # Add timeline items
            for i, item in enumerate(roadmap_items):
                x = timeline_start_x + (i * item_width) + (item_width / 2)
                
                # Timeline point
                point = shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(x - 0.1), Inches(timeline_y - 0.1),
                    Inches(0.2), Inches(0.2)
                )
                
                if hasattr(point, 'fill'):
                    point.fill.solid()
                    point.fill.fore_color.rgb = self._get_status_color(item.status)
                
                # Item label above timeline
                label_box = shapes.add_textbox(
                    Inches(x - 0.8), Inches(timeline_y - 1.5),
                    Inches(1.6), Inches(1.2)
                )
                
                text_frame = label_box.text_frame
                text_frame.clear()
                
                # Title
                p = text_frame.paragraphs[0]
                p.text = item.title
                p.font.bold = True
                p.font.size = Pt(12)
                p.font.color.rgb = self._get_theme_color('text1')
                p.alignment = 1  # Center
                
                # Description
                if item.description:
                    p = text_frame.add_paragraph()
                    p.text = item.description[:50] + ('...' if len(item.description) > 50 else '')
                    p.font.size = Pt(10)
                    p.font.color.rgb = self._get_theme_color('text2')
                    p.alignment = 1  # Center
                    p.space_before = Pt(3)
                
                # Status below timeline
                status_box = shapes.add_textbox(
                    Inches(x - 0.5), Inches(timeline_y + 0.3),
                    Inches(1), Inches(0.4)
                )
                
                status_frame = status_box.text_frame
                status_frame.clear()
                
                status_p = status_frame.paragraphs[0]
                status_p.text = item.status.replace('_', ' ').title()
                status_p.font.size = Pt(9)
                status_p.font.color.rgb = self._get_status_color(item.status)
                status_p.alignment = 1  # Center
                
        except Exception:
            # Fallback to simple text list
            self._create_simple_roadmap_text(slide, roadmap_items)
    
    def _create_vertical_timeline(self, slide: object, roadmap_items: List[RoadmapItem]) -> None:
        """Create vertical timeline layout for many items."""
        try:
            shapes = getattr(slide, 'shapes', None)
            if not shapes or not hasattr(shapes, 'add_shape'):
                return
            
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.util import Inches, Pt
            
            # Timeline parameters
            timeline_x = 2
            timeline_height = 5
            timeline_start_y = 1.5
            item_height = timeline_height / len(roadmap_items)
            
            # Draw timeline base line
            line = shapes.add_connector(
                1,  # Straight line
                Inches(timeline_x), Inches(timeline_start_y),
                Inches(timeline_x), Inches(timeline_start_y + timeline_height)
            )
            
            if hasattr(line, 'line'):
                line.line.color.rgb = self._get_theme_color('accent1')
                line.line.width = Inches(0.05)
            
            # Add timeline items
            for i, item in enumerate(roadmap_items):
                y = timeline_start_y + (i * item_height) + (item_height / 2)
                
                # Timeline point
                point = shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(timeline_x - 0.1), Inches(y - 0.1),
                    Inches(0.2), Inches(0.2)
                )
                
                if hasattr(point, 'fill'):
                    point.fill.solid()
                    point.fill.fore_color.rgb = self._get_status_color(item.status)
                
                # Item content to the right
                content_box = shapes.add_textbox(
                    Inches(timeline_x + 0.3), Inches(y - 0.3),
                    Inches(6), Inches(0.6)
                )
                
                text_frame = content_box.text_frame
                text_frame.clear()
                
                # Title and description
                p = text_frame.paragraphs[0]
                p.text = f"{item.title} - {item.status.replace('_', ' ').title()}"
                p.font.bold = True
                p.font.size = Pt(12)
                p.font.color.rgb = self._get_theme_color('text1')
                
                if item.description:
                    p = text_frame.add_paragraph()
                    p.text = item.description
                    p.font.size = Pt(10)
                    p.font.color.rgb = self._get_theme_color('text2')
                
        except Exception:
            # Fallback to simple text list
            self._create_simple_roadmap_text(slide, roadmap_items)
    
    def _create_timeline_section(self, slide: object, roadmap_items: List[RoadmapItem], 
                               x: float, y: float, width: float, height: float) -> None:
        """Create a timeline section within specified bounds."""
        # Use horizontal layout for section
        try:
            shapes = getattr(slide, 'shapes', None)
            if not shapes:
                return
            
            # Section title
            if hasattr(shapes, 'add_textbox'):
                from pptx.util import Inches, Pt
                
                title_box = shapes.add_textbox(
                    Inches(x), Inches(y),
                    Inches(width), Inches(0.5)
                )
                
                title_frame = title_box.text_frame
                title_p = title_frame.paragraphs[0]
                title_p.text = "Timeline"
                title_p.font.bold = True
                title_p.font.size = Pt(16)
                title_p.font.color.rgb = self._get_theme_color('accent1')
                
                # Create mini timeline
                self._create_mini_timeline(slide, roadmap_items, x, y + 0.7, width, height - 0.7)
                
        except Exception:
            pass
    
    def _create_milestones_section(self, slide: object, milestones: List[Milestone], 
                                 x: float, y: float, width: float, height: float) -> None:
        """Create milestones section."""
        try:
            shapes = getattr(slide, 'shapes', None)
            if not shapes or not hasattr(shapes, 'add_textbox'):
                return
            
            from pptx.util import Inches, Pt
            
            # Section title
            title_box = shapes.add_textbox(
                Inches(x), Inches(y),
                Inches(width), Inches(0.5)
            )
            
            title_frame = title_box.text_frame
            title_p = title_frame.paragraphs[0]
            title_p.text = "Key Milestones"
            title_p.font.bold = True
            title_p.font.size = Pt(16)
            title_p.font.color.rgb = self._get_theme_color('accent1')
            
            # Milestones list
            milestone_height = (height - 0.7) / len(milestones) if milestones else 0.5
            
            for i, milestone in enumerate(milestones):
                milestone_y = y + 0.7 + (i * milestone_height)
                
                milestone_box = shapes.add_textbox(
                    Inches(x + 0.5), Inches(milestone_y),
                    Inches(width - 0.5), Inches(milestone_height)
                )
                
                milestone_frame = milestone_box.text_frame
                milestone_frame.clear()
                
                # Milestone title and date
                p = milestone_frame.paragraphs[0]
                date_str = milestone.date.strftime('%m/%d/%Y') if milestone.date else 'TBD'
                p.text = f"• {milestone.title} ({date_str})"
                p.font.size = Pt(12)
                p.font.color.rgb = self._get_theme_color('text1')
                
                # Description if available
                if milestone.description:
                    p = milestone_frame.add_paragraph()
                    p.text = f"  {milestone.description}"
                    p.font.size = Pt(10)
                    p.font.color.rgb = self._get_theme_color('text2')
                
        except Exception:
            pass
    
    def _create_mini_timeline(self, slide: object, roadmap_items: List[RoadmapItem], 
                            x: float, y: float, width: float, height: float) -> None:
        """Create a condensed timeline view."""
        try:
            shapes = getattr(slide, 'shapes', None)
            if not shapes or not roadmap_items:
                return
            
            from pptx.util import Inches, Pt
            
            # Simple progress bar style timeline
            bar_width = width - 1
            bar_height = 0.3
            bar_y = y + (height / 2) - (bar_height / 2)
            
            # Background bar
            if hasattr(shapes, 'add_shape'):
                from pptx.enum.shapes import MSO_SHAPE
                
                bg_bar = shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(x + 0.5), Inches(bar_y),
                    Inches(bar_width), Inches(bar_height)
                )
                
                if hasattr(bg_bar, 'fill'):
                    bg_bar.fill.solid()
                    bg_bar.fill.fore_color.rgb = self._get_theme_color('accent2')
                
                # Progress indicators
                section_width = bar_width / len(roadmap_items)
                
                for i, item in enumerate(roadmap_items):
                    section_x = x + 0.5 + (i * section_width)
                    
                    # Status indicator
                    indicator = shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(section_x), Inches(bar_y),
                        Inches(section_width * 0.8), Inches(bar_height)
                    )
                    
                    if hasattr(indicator, 'fill'):
                        indicator.fill.solid()
                        indicator.fill.fore_color.rgb = self._get_status_color(item.status)
                    
                    # Item label below
                    label_box = shapes.add_textbox(
                        Inches(section_x), Inches(bar_y + bar_height + 0.1),
                        Inches(section_width), Inches(0.4)
                    )
                    
                    label_frame = label_box.text_frame
                    label_p = label_frame.paragraphs[0]
                    label_p.text = item.title
                    label_p.font.size = Pt(9)
                    label_p.font.color.rgb = self._get_theme_color('text2')
                    label_p.alignment = 1  # Center
                    
        except Exception:
            pass
    
    def _create_default_layout(self, slide: object, slide_data: Dict[str, Any]) -> None:
        """Create default layout when no specific data is available."""
        try:
            content = slide_data.get('content', slide_data.get('points', []))
            if content:
                self._add_bullet_points(slide, content)
        except Exception:
            pass
    
    def _add_bullet_points(self, slide: object, content: List[str]) -> None:
        """Add bullet points to slide."""
        try:
            if not content:
                return
            
            shapes = getattr(slide, 'shapes', None)
            if shapes and hasattr(shapes, 'add_textbox'):
                from pptx.util import Inches
                
                textbox = shapes.add_textbox(
                    Inches(1), Inches(2),
                    Inches(8), Inches(4)
                )
                
                text_frame = textbox.text_frame
                text_frame.clear()
                
                for i, point in enumerate(content):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = f"• {point}"
                    p.level = 0
        except Exception:
            pass
    
    def _create_simple_roadmap_text(self, slide: object, roadmap_items: List[RoadmapItem]) -> None:
        """Create simple text representation of roadmap."""
        try:
            shapes = getattr(slide, 'shapes', None)
            if shapes and hasattr(shapes, 'add_textbox'):
                from pptx.util import Inches
                
                content = []
                for item in roadmap_items:
                    item_text = f"• {item.title}"
                    if item.status:
                        item_text += f" ({item.status.replace('_', ' ').title()})"
                    if item.description:
                        item_text += f": {item.description}"
                    content.append(item_text)
                
                textbox = shapes.add_textbox(
                    Inches(1), Inches(2),
                    Inches(8), Inches(4)
                )
                textbox.text = "\n".join(content)
        except Exception:
            pass
    
    def _get_theme_color(self, color_name: str) -> object:
        """Get theme color."""
        try:
            from pptx.dml.color import RGBColor
            
            colors = {
                'accent1': RGBColor(68, 114, 196),
                'accent2': RGBColor(112, 173, 71),
                'text1': RGBColor(68, 68, 68),
                'text2': RGBColor(89, 89, 89),
            }
            return colors.get(color_name, RGBColor(68, 68, 68))
        except ImportError:
            return None
    
    def _get_status_color(self, status: str) -> object:
        """Get color for status indicators."""
        try:
            from pptx.dml.color import RGBColor
            
            colors = {
                'completed': RGBColor(34, 139, 34),      # Green
                'in_progress': RGBColor(255, 165, 0),    # Orange
                'planned': RGBColor(68, 114, 196),       # Blue
                'delayed': RGBColor(220, 20, 60),        # Red
            }
            return colors.get(status, RGBColor(68, 114, 196))
        except ImportError:
            return None
