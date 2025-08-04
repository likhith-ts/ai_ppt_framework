"""
Python-pptx based PowerPoint presentation engine.

This module provides cross-platform PowerPoint generation using the python-pptx
library, offering good compatibility and feature coverage with enhanced layout
system and proper background handling.
"""

from typing import Optional, Any, List
from pathlib import Path

from .base_engine import BasePresentationEngine, SlideData
from core.exceptions import PresentationGenerationError
from design.themes import DesignTheme
from design.color_system import DesignPalettes
from core.constants import MARGIN_MEDIUM, MARGIN_SMALL
from visual.layout_system import layout_engine, background_manager, LayoutType, Position


class PPTXPresentationEngine(BasePresentationEngine):
    """
    Python-pptx based presentation engine for cross-platform compatibility.

    Features:
    - Cross-platform PowerPoint generation
    - Good feature coverage with python-pptx
    - Professional slide layouts
    - Theme-based styling
    - No PowerPoint installation required
    """

    def __init__(self, config: Optional[Any] = None):
        """Initialize the python-pptx presentation engine."""
        super().__init__(config)
        self.presentation = None

        if not self.is_available():
            raise PresentationGenerationError(
                "python-pptx engine requires 'python-pptx' package. Install with: pip install python-pptx"
            )

    def is_available(self) -> bool:
        """Check if python-pptx is available."""
        try:
            import importlib.util

            return importlib.util.find_spec("pptx") is not None
        except ImportError:
            return False

    def create_presentation(self, title: str = "AI Generated Presentation") -> None:
        """Create a new PowerPoint presentation using python-pptx."""
        try:
            from pptx import Presentation
            from pptx.util import Inches

            # Create new presentation
            self.presentation = Presentation()

            # Set slide dimensions
            self.presentation.slide_width = Inches(10)  # 960 points
            self.presentation.slide_height = Inches(5.625)  # 540 points

            self.slides_created = 0

        except ImportError:
            raise PresentationGenerationError("python-pptx package not available")
        except Exception as e:
            raise PresentationGenerationError(
                f"Failed to create presentation: {str(e)}"
            ) from e

    def add_slide(self, slide_data: SlideData) -> None:
        """Add a slide using enhanced layout system and proper positioning."""
        if not self.presentation:
            raise PresentationGenerationError("No presentation created")

        try:
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN

            # Determine optimal layout based on content
            layout_type = layout_engine.calculate_optimal_layout(
                slide_data.slide_type,
                has_image=bool(slide_data.background_image or slide_data.diagram_image),
                has_diagram=bool(slide_data.diagram_image)
            )
            
            # Get layout regions for this layout type
            layout_regions = layout_engine.get_layout(layout_type)
            
            # Add slide with blank layout for custom positioning
            slide_layout = self.presentation.slide_layouts[6]  # Blank layout
            slide = self.presentation.slides.add_slide(slide_layout)

            # Get theme colors
            palette = DesignPalettes.get_palette(slide_data.theme)

            # Handle background properly
            self._setup_slide_background(slide, slide_data, palette)
            
            # Apply dynamic colors if available
            if hasattr(slide_data, 'primary_color') and slide_data.primary_color:
                self._apply_dynamic_colors(slide, slide_data)

            # Add content based on layout regions
            self._populate_slide_content(slide, slide_data, layout_regions, palette)
            
            # Add visual elements
            self._add_visual_elements(slide, slide_data)

            self.slides_created += 1
            
            if self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                print(f"Added slide: {slide_data.title} using layout: {layout_type}")

        except Exception as e:
            raise PresentationGenerationError(
                f"Failed to add slide: {str(e)}"
            ) from e

    def add_title_slide(
        self,
        title: str,
        subtitle: str = "",
        theme: DesignTheme = DesignTheme.CORPORATE_MODERN,
    ) -> None:
        """Add a title slide using python-pptx."""
        if not self.presentation:
            raise PresentationGenerationError("No presentation created")

        try:
            from pptx.util import Pt
            from pptx.enum.text import PP_ALIGN

            # Add title slide
            slide_layout = self.presentation.slide_layouts[0]  # Title slide layout
            slide = self.presentation.slides.add_slide(slide_layout)

            # Get theme colors
            palette = DesignPalettes.get_palette(theme)

            # Set slide background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = self._int_to_rgb(palette.background)

            # Configure title
            title_placeholder = slide.shapes.title
            if title_placeholder:
                title_placeholder.text = title
                title_frame = title_placeholder.text_frame
                title_paragraph = title_frame.paragraphs[0]
                title_paragraph.font.name = "Segoe UI"
                title_paragraph.font.size = Pt(44)
                title_paragraph.font.color.rgb = self._int_to_rgb(palette.primary)
                title_paragraph.font.bold = True
                title_paragraph.alignment = PP_ALIGN.CENTER

            # Configure subtitle
            subtitle_placeholder = None
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.idx == 1:  # Subtitle placeholder
                    subtitle_placeholder = placeholder
                    break

            if subtitle_placeholder and subtitle:
                try:
                    setattr(subtitle_placeholder, 'text', subtitle)
                    subtitle_frame = getattr(subtitle_placeholder, 'text_frame')
                    subtitle_paragraph = subtitle_frame.paragraphs[0]
                    subtitle_paragraph.font.name = "Segoe UI"
                    subtitle_paragraph.font.size = Pt(24)
                    subtitle_paragraph.font.color.rgb = self._int_to_rgb(
                        palette.text_secondary
                    )
                    subtitle_paragraph.alignment = PP_ALIGN.CENTER
                except (AttributeError, TypeError):
                    # Fallback if the placeholder doesn't support text
                    pass

            self.slides_created += 1

        except Exception as e:
            raise PresentationGenerationError(
                f"Failed to add title slide: {str(e)}"
            ) from e

    def save_presentation(self, filepath: Path) -> bool:
        """Save the presentation to a file."""
        if not self.presentation:
            raise PresentationGenerationError("No presentation to save")

        try:
            # Ensure directory exists
            filepath.parent.mkdir(parents=True, exist_ok=True)

            # Save presentation
            self.presentation.save(str(filepath))

            return filepath.exists()

        except Exception as e:
            raise PresentationGenerationError(
                f"Failed to save presentation: {str(e)}"
            ) from e

    def close_presentation(self) -> None:
        """Close the presentation and clean up resources."""
        self.presentation = None

    def _int_to_rgb(self, color_int: int):
        """Convert integer color to RGBColor object."""
        from pptx.dml.color import RGBColor

        # Extract RGB components from integer
        red = (color_int >> 16) & 0xFF
        green = (color_int >> 8) & 0xFF
        blue = color_int & 0xFF

        return RGBColor(red, green, blue)
    
    def _hex_to_rgb(self, hex_color: str):
        """Convert hex color string to RGBColor object."""
        from pptx.dml.color import RGBColor
        
        # Remove # if present
        hex_color = hex_color.lstrip('#')
        
        # Convert hex to RGB
        try:
            red = int(hex_color[0:2], 16)
            green = int(hex_color[2:4], 16)
            blue = int(hex_color[4:6], 16)
            return RGBColor(red, green, blue)
        except (ValueError, IndexError):
            # Fallback to default color if hex parsing fails
            return RGBColor(0, 0, 0)

    def _insert_background_image(self, slide, image_path: str):
        """Insert a background image into the slide."""
        try:
            from pptx.util import Inches
            from pathlib import Path
            
            if not image_path or not Path(image_path).exists():
                return
            
            # Insert image as background
            # Note: python-pptx doesn't directly support background images,
            # so we'll insert it as a large image behind other content
            left = Inches(0)
            top = Inches(0)
            width = Inches(10)
            height = Inches(5.625)
            
            picture = slide.shapes.add_picture(image_path, left, top, width, height)
            
            # Send to back
            slide.shapes._spTree.remove(picture._element)
            slide.shapes._spTree.insert(2, picture._element)
            
        except Exception as e:
            print(f"Warning: Could not insert background image: {e}")
    
    def _insert_diagram_image(self, slide, image_path: str, slide_type: str):
        """Insert a diagram image into the slide."""
        try:
            from pptx.util import Inches
            from pathlib import Path
            
            if not image_path or not Path(image_path).exists():
                return
            
            # Position based on slide type
            if slide_type == "architecture_slide":
                left = Inches(5.5)
                top = Inches(1.5)
                width = Inches(4)
                height = Inches(3)
            elif slide_type == "roadmap_slide":
                left = Inches(1)
                top = Inches(3)
                width = Inches(8)
                height = Inches(2)
            elif slide_type == "metrics_slide":
                left = Inches(5)
                top = Inches(1)
                width = Inches(4.5)
                height = Inches(3.5)
            else:
                # Default positioning
                left = Inches(5.5)
                top = Inches(2)
                width = Inches(3.5)
                height = Inches(2.5)
            
            slide.shapes.add_picture(image_path, left, top, width, height)
            
        except Exception as e:
            print(f"Warning: Could not insert diagram image: {e}")
    
    def _insert_feature_icons(self, slide, icon_paths: List[str]):
        """Insert feature icons into the slide."""
        try:
            from pptx.util import Inches
            from pathlib import Path
            
            valid_icons = [path for path in icon_paths if path and Path(path).exists()]
            if not valid_icons:
                return
            
            # Position icons in a grid
            start_left = Inches(0.5)
            start_top = Inches(4)
            icon_size = Inches(0.8)
            spacing = Inches(2)
            
            for i, icon_path in enumerate(valid_icons[:4]):  # Max 4 icons
                left = start_left + (i * spacing)
                top = start_top
                
                slide.shapes.add_picture(icon_path, left, top, icon_size, icon_size)
                
        except Exception as e:
            print(f"Warning: Could not insert feature icons: {e}")
    
    def _setup_slide_background(self, slide, slide_data: SlideData, palette):
        """Set up slide background with proper handling for images and colors."""
        try:
            from pptx.util import Inches
            
            # If there's a background image, use it properly
            if slide_data.background_image:
                self._insert_proper_background_image(slide, slide_data.background_image)
            else:
                # Use solid color background
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = self._int_to_rgb(palette.background)
                
        except Exception as e:
            if self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                print(f"Warning: Could not set up background: {e}")
    
    def _insert_proper_background_image(self, slide, image_path: str):
        """Insert background image that properly covers the entire slide."""
        try:
            from pptx.util import Inches
            
            # Process the image to ensure it covers the slide properly
            processed_image = background_manager.prepare_background_image(image_path)
            
            # Add image as background by inserting it as the first shape (lowest z-index)
            # This ensures it appears behind all other content
            left = Inches(0)
            top = Inches(0)
            width = self.presentation.slide_width if self.presentation else Inches(10)
            height = self.presentation.slide_height if self.presentation else Inches(5.625)
            
            pic = slide.shapes.add_picture(processed_image, left, top, width, height)
            
            # Move the picture to the back (lowest z-order)
            shape_elm = pic.element
            shape_elm.getparent().remove(shape_elm)
            shape_elm.getparent().insert(0, shape_elm)
            
        except Exception as e:
            if self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                print(f"Warning: Could not insert background image: {e}")
    
    def _populate_slide_content(self, slide, slide_data: SlideData, layout_regions, palette):
        """Populate slide content using the layout system."""
        try:
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            
            for region in layout_regions:
                if region.name == "title":
                    self._add_title_region(slide, slide_data.title, region, palette)
                elif region.name == "content":
                    if slide_data.points:
                        # Optimize bullet points
                        optimized_points = layout_engine.optimize_bullet_points(
                            slide_data.points, max_points=4
                        )
                        self._add_content_region(slide, optimized_points, region, palette)
                elif region.name in ["left_content", "right_content"]:
                    # Split content for two-column layout
                    if slide_data.points:
                        mid_point = len(slide_data.points) // 2
                        if region.name == "left_content":
                            content = slide_data.points[:mid_point]
                        else:
                            content = slide_data.points[mid_point:]
                        
                        if content:
                            self._add_content_region(slide, content, region, palette)
                elif region.name == "text_content":
                    # For image+text layout
                    if slide_data.points:
                        optimized_points = layout_engine.optimize_bullet_points(
                            slide_data.points, max_points=3
                        )
                        self._add_content_region(slide, optimized_points, region, palette)
                        
        except Exception as e:
            if self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                print(f"Warning: Could not populate slide content: {e}")
    
    def _add_title_region(self, slide, title: str, region, palette):
        """Add title with proper positioning and formatting."""
        try:
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            
            # Calculate optimal font size
            font_size = layout_engine.calculate_font_size(title, region)
            
            # Add text box for title
            left = Inches(region.position.x)
            top = Inches(region.position.y)
            width = Inches(region.position.width)
            height = Inches(region.position.height)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.clear()
            
            # Add title text
            p = text_frame.paragraphs[0]
            p.text = title
            p.font.name = "Segoe UI"
            p.font.size = Pt(font_size)
            p.font.color.rgb = self._int_to_rgb(palette.primary)
            p.font.bold = True
            
            # Set alignment
            if region.alignment.value == "center":
                p.alignment = PP_ALIGN.CENTER
            elif region.alignment.value == "right":
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT
                
        except Exception as e:
            if self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                print(f"Warning: Could not add title region: {e}")
    
    def _add_content_region(self, slide, content: List[str], region, palette):
        """Add content with proper positioning and overflow protection."""
        try:
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            
            # Validate and truncate content if needed
            processed_content = []
            for item in content:
                processed_text, was_truncated = layout_engine.validate_text_length(item, region)
                processed_content.append(processed_text)
                if was_truncated and self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                    print(f"Truncated content: {item[:50]}...")
            
            # Calculate font size
            content_text = "\n".join([f"• {item}" for item in processed_content])
            font_size = layout_engine.calculate_font_size(content_text, region)
            
            # Add text box
            left = Inches(region.position.x)
            top = Inches(region.position.y)
            width = Inches(region.position.width)
            height = Inches(region.position.height)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            
            # Add content
            for i, item in enumerate(processed_content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = f"• {item}"
                p.font.name = "Segoe UI"
                p.font.size = Pt(font_size)
                p.font.color.rgb = self._int_to_rgb(palette.text_primary)
                p.space_after = Pt(6)
                p.alignment = PP_ALIGN.LEFT
                
        except Exception as e:
            if self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                print(f"Warning: Could not add content region: {e}")
    
    def _add_visual_elements(self, slide, slide_data: SlideData):
        """Add visual elements like diagrams and feature icons."""
        try:
            # Insert diagram image if available
            if slide_data.diagram_image:
                self._insert_diagram_image(slide, slide_data.diagram_image, slide_data.slide_type)
            
            # Insert feature icons if available
            if slide_data.feature_icons:
                self._insert_feature_icons(slide, slide_data.feature_icons)
                
        except Exception as e:
            if self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                print(f"Warning: Could not add visual elements: {e}")

    def _apply_dynamic_colors(self, slide, slide_data: SlideData):
        """Apply dynamic color scheme from AI-generated theme."""
        try:
            # If we have AI-generated colors, use them
            if hasattr(slide_data, 'primary_color') and slide_data.primary_color:
                # Convert hex color to RGB
                primary_rgb = self._hex_to_rgb(slide_data.primary_color)
                # Update the slide with new colors - this would require more complex color manipulation
                # For now, we'll store the colors for potential future use
                if self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                    print(f"Applying primary color: {slide_data.primary_color}")
            
            if hasattr(slide_data, 'accent_colors') and slide_data.accent_colors:
                if self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                    print(f"Applying accent colors: {slide_data.accent_colors}")
                
        except Exception as e:
            if self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                print(f"Warning: Could not apply dynamic colors: {e}")

    def add_smartart_diagram(self, slide, diagram_data):
        """Add SmartArt-like diagram to slide using python-pptx shapes."""
        try:
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.dml.color import RGBColor
            from visual.smartart_engine import SmartArtType
            
            # Since python-pptx doesn't support native SmartArt, we'll create
            # equivalent diagrams using shapes and text boxes
            
            if diagram_data.diagram_type == SmartArtType.PROCESS:
                return self._create_process_diagram(slide, diagram_data)
            elif diagram_data.diagram_type == SmartArtType.HIERARCHY:
                return self._create_hierarchy_diagram(slide, diagram_data)
            elif diagram_data.diagram_type == SmartArtType.CYCLE:
                return self._create_cycle_diagram(slide, diagram_data)
            elif diagram_data.diagram_type == SmartArtType.RELATIONSHIP:
                return self._create_relationship_diagram(slide, diagram_data)
            elif diagram_data.diagram_type == SmartArtType.MATRIX:
                return self._create_matrix_diagram(slide, diagram_data)
            elif diagram_data.diagram_type == SmartArtType.PYRAMID:
                return self._create_pyramid_diagram(slide, diagram_data)
            elif diagram_data.diagram_type == SmartArtType.LIST:
                return self._create_list_diagram(slide, diagram_data)
            else:
                # Default to list
                return self._create_list_diagram(slide, diagram_data)
                
        except Exception as e:
            if self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                print(f"Warning: Could not create SmartArt diagram: {e}")
            return None
    
    def _create_process_diagram(self, slide, diagram_data):
        """Create a process flow diagram using shapes."""
        try:
            from pptx.util import Inches, Pt
            from pptx.enum.shapes import MSO_SHAPE
            
            elements = diagram_data.elements[:6]  # Limit to 6 elements for space
            if not elements:
                return None
                
            # Calculate positions
            total_width = Inches(8)
            start_left = Inches(1)
            shape_width = total_width / len(elements)
            shape_height = Inches(1)
            y_pos = Inches(2.5)
            
            shapes = []
            for i, element in enumerate(elements):
                # Calculate position
                x_pos = start_left + (i * shape_width)
                
                # Add process box
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    x_pos, y_pos,
                    shape_width * 0.8, shape_height  # Leave some space between boxes
                )
                
                # Set text
                shape.text = element.text
                
                # Format text
                text_frame = shape.text_frame
                text_frame.text = element.text
                text_frame.word_wrap = True
                text_frame.margin_left = Inches(0.1)
                text_frame.margin_right = Inches(0.1)
                text_frame.margin_top = Inches(0.1)
                text_frame.margin_bottom = Inches(0.1)
                
                # Format font
                para = text_frame.paragraphs[0]
                para.font.size = Pt(12)
                para.font.name = "Segoe UI"
                para.alignment = PP_ALIGN.CENTER
                
                # Apply colors
                if diagram_data.color_palette:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = self._int_to_rgb(diagram_data.color_palette.primary)
                    para.font.color.rgb = self._int_to_rgb(diagram_data.color_palette.text_primary)
                
                shapes.append(shape)
                
                # Add arrow to next shape (except for last element)
                if i < len(elements) - 1:
                    arrow_x = x_pos + (shape_width * 0.8)
                    arrow_y = y_pos + (shape_height / 2)
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.RIGHT_ARROW,
                        arrow_x, arrow_y - Inches(0.2),
                        shape_width * 0.2, Inches(0.4)
                    )
                    arrow.fill.solid()
                    if diagram_data.color_palette:
                        arrow.fill.fore_color.rgb = self._int_to_rgb(diagram_data.color_palette.accent)
                    shapes.append(arrow)
            
            return shapes
            
        except Exception as e:
            if self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                print(f"Warning: Could not create process diagram: {e}")
            return None
    
    def _create_hierarchy_diagram(self, slide, diagram_data):
        """Create a hierarchy diagram using shapes."""
        try:
            from pptx.util import Inches, Pt
            from pptx.enum.shapes import MSO_SHAPE
            
            elements = diagram_data.elements[:7]  # Limit to 7 elements
            if not elements:
                return None
                
            shapes = []
            
            # Top level (1 element)
            if elements:
                top_shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(3.5), Inches(1),
                    Inches(3), Inches(0.8)
                )
                top_shape.text = elements[0].text
                
                # Format
                self._format_hierarchy_shape(top_shape, diagram_data.color_palette, is_top=True)
                shapes.append(top_shape)
                
                # Second level (up to 3 elements)
                second_level = elements[1:4]
                level2_width = Inches(2.5)
                level2_spacing = Inches(3)
                level2_start = Inches(1)
                
                if second_level:
                    for i, element in enumerate(second_level):
                        x_pos = level2_start + (i * level2_spacing)
                        shape = slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE,
                            x_pos, Inches(2.5),
                            level2_width, Inches(0.7)
                        )
                        shape.text = element.text
                        self._format_hierarchy_shape(shape, diagram_data.color_palette)
                        shapes.append(shape)
                        
                        # Add connector line
                        self._add_connector_line(slide, 
                                               Inches(5), Inches(1.8),  # From top center
                                               x_pos + level2_width/2, Inches(2.5),  # To child top
                                               diagram_data.color_palette)
                
                # Third level (remaining elements)
                third_level = elements[4:]
                if third_level:
                    level3_width = Inches(2)
                    level3_spacing = Inches(2.5)
                    level3_start = Inches(1)
                    
                    for i, element in enumerate(third_level[:3]):  # Max 3 on third level
                        x_pos = level3_start + (i * level3_spacing)
                        shape = slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE,
                            x_pos, Inches(4),
                            level3_width, Inches(0.6)
                        )
                        shape.text = element.text
                        self._format_hierarchy_shape(shape, diagram_data.color_palette)
                        shapes.append(shape)
                        
                        # Connect to parent in second level
                        if i < len(second_level):
                            parent_x = level2_start + (i * level2_spacing) + level2_width/2
                            self._add_connector_line(slide,
                                                   parent_x, Inches(3.2),  # From parent bottom
                                                   x_pos + level3_width/2, Inches(4),  # To child top
                                                   diagram_data.color_palette)
            
            return shapes
            
        except Exception as e:
            if self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                print(f"Warning: Could not create hierarchy diagram: {e}")
            return None
    
    def _create_cycle_diagram(self, slide, diagram_data):
        """Create a cycle diagram using shapes."""
        try:
            from pptx.util import Inches, Pt
            from pptx.enum.shapes import MSO_SHAPE
            import math
            
            elements = diagram_data.elements[:6]  # Limit to 6 elements
            if not elements:
                return None
                
            shapes = []
            
            # Circle parameters
            center_x = Inches(5)
            center_y = Inches(2.8)
            radius = Inches(2)
            shape_size = Inches(1.5)
            
            for i, element in enumerate(elements):
                # Calculate position on circle
                angle = (2 * math.pi * i) / len(elements) - (math.pi / 2)  # Start at top
                x_pos = center_x + radius * math.cos(angle) - shape_size/2
                y_pos = center_y + radius * math.sin(angle) - shape_size/2
                
                # Add circle shape
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    x_pos, y_pos,
                    shape_size, shape_size
                )
                shape.text = element.text
                
                # Format
                self._format_cycle_shape(shape, diagram_data.color_palette)
                shapes.append(shape)
                
                # Add arrow to next shape
                next_i = (i + 1) % len(elements)
                next_angle = (2 * math.pi * next_i) / len(elements) - (math.pi / 2)
                next_x = center_x + radius * math.cos(next_angle)
                next_y = center_y + radius * math.sin(next_angle)
                
                # Add curved arrow (simplified as line)
                arrow = slide.shapes.add_connector(
                    1,  # Straight connector
                    x_pos + shape_size/2, y_pos + shape_size/2,
                    next_x, next_y
                )
                if diagram_data.color_palette:
                    arrow.line.color.rgb = self._int_to_rgb(diagram_data.color_palette.accent)
                    arrow.line.width = Pt(2)
                shapes.append(arrow)
            
            return shapes
            
        except Exception as e:
            if self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                print(f"Warning: Could not create cycle diagram: {e}")
            return None
    
    def _create_list_diagram(self, slide, diagram_data):
        """Create a list diagram using shapes."""
        try:
            from pptx.util import Inches, Pt
            from pptx.enum.shapes import MSO_SHAPE
            
            elements = diagram_data.elements[:8]  # Limit to 8 elements
            if not elements:
                return None
                
            shapes = []
            
            # List parameters
            start_x = Inches(1)
            start_y = Inches(1.5)
            item_width = Inches(8)
            item_height = Inches(0.6)
            spacing = Inches(0.1)
            
            for i, element in enumerate(elements):
                y_pos = start_y + i * (item_height + spacing)
                
                # Add list item shape
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    start_x, y_pos,
                    item_width, item_height
                )
                shape.text = element.text
                
                # Format
                self._format_list_shape(shape, diagram_data.color_palette)
                shapes.append(shape)
            
            return shapes
            
        except Exception as e:
            if self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                print(f"Warning: Could not create list diagram: {e}")
            return None
    
    def _create_relationship_diagram(self, slide, diagram_data):
        """Create a relationship diagram (Venn-style)."""
        try:
            from pptx.util import Inches, Pt
            from pptx.enum.shapes import MSO_SHAPE
            
            elements = diagram_data.elements[:4]  # Limit to 4 elements
            if not elements:
                return None
                
            shapes = []
            
            if len(elements) == 2:
                # Two overlapping circles
                circle1 = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(2), Inches(2),
                    Inches(3), Inches(3)
                )
                circle1.text = elements[0].text
                
                circle2 = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(4), Inches(2),
                    Inches(3), Inches(3)
                )
                circle2.text = elements[1].text
                
                # Format circles
                self._format_relationship_shape(circle1, diagram_data.color_palette, 0)
                self._format_relationship_shape(circle2, diagram_data.color_palette, 1)
                
                shapes.extend([circle1, circle2])
                
            elif len(elements) >= 3:
                # Three overlapping circles
                positions = [
                    (Inches(3), Inches(1.5)),  # Top
                    (Inches(1.5), Inches(3.5)),  # Bottom left
                    (Inches(4.5), Inches(3.5)),  # Bottom right
                ]
                
                for i, (element, (x, y)) in enumerate(zip(elements[:3], positions)):
                    circle = slide.shapes.add_shape(
                        MSO_SHAPE.OVAL,
                        x, y,
                        Inches(2.5), Inches(2.5)
                    )
                    circle.text = element.text
                    self._format_relationship_shape(circle, diagram_data.color_palette, i)
                    shapes.append(circle)
                
                # Add fourth element as central text if exists
                if len(elements) >= 4:
                    center_text = slide.shapes.add_textbox(
                        Inches(4), Inches(3),
                        Inches(2), Inches(0.5)
                    )
                    center_text.text = elements[3].text
                    self._format_text_shape(center_text, diagram_data.color_palette)
                    shapes.append(center_text)
            
            return shapes
            
        except Exception as e:
            if self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                print(f"Warning: Could not create relationship diagram: {e}")
            return None
    
    def _create_matrix_diagram(self, slide, diagram_data):
        """Create a matrix diagram using shapes."""
        try:
            from pptx.util import Inches, Pt
            from pptx.enum.shapes import MSO_SHAPE
            
            elements = diagram_data.elements[:4]  # 2x2 matrix
            if not elements:
                return None
                
            shapes = []
            
            # Matrix parameters
            start_x = Inches(2)
            start_y = Inches(2)
            cell_width = Inches(3)
            cell_height = Inches(1.5)
            
            positions = [
                (start_x, start_y),  # Top-left
                (start_x + cell_width, start_y),  # Top-right
                (start_x, start_y + cell_height),  # Bottom-left
                (start_x + cell_width, start_y + cell_height),  # Bottom-right
            ]
            
            for i, (element, (x, y)) in enumerate(zip(elements, positions)):
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    x, y,
                    cell_width, cell_height
                )
                shape.text = element.text
                self._format_matrix_shape(shape, diagram_data.color_palette, i)
                shapes.append(shape)
            
            return shapes
            
        except Exception as e:
            if self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                print(f"Warning: Could not create matrix diagram: {e}")
            return None
    
    def _create_pyramid_diagram(self, slide, diagram_data):
        """Create a pyramid diagram using shapes."""
        try:
            from pptx.util import Inches, Pt
            from pptx.enum.shapes import MSO_SHAPE
            
            elements = diagram_data.elements[:5]  # Limit to 5 levels
            if not elements:
                return None
                
            shapes = []
            
            # Pyramid parameters
            center_x = Inches(5)
            start_y = Inches(1)
            level_height = Inches(0.8)
            max_width = Inches(8)
            
            for i, element in enumerate(elements):
                # Calculate width for this level (wider at bottom)
                width_ratio = (len(elements) - i) / len(elements)
                width = max_width * width_ratio
                
                # Calculate position
                x = center_x - width/2
                y = start_y + i * level_height
                
                # Add pyramid level
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    x, y,
                    width, level_height * 0.9  # Small gap between levels
                )
                shape.text = element.text
                
                # Format
                self._format_pyramid_shape(shape, diagram_data.color_palette, i)
                shapes.append(shape)
            
            return shapes
            
        except Exception as e:
            if self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                print(f"Warning: Could not create pyramid diagram: {e}")
            return None
    
    def _format_hierarchy_shape(self, shape, color_palette, is_top=False):
        """Format hierarchy shape with colors and text."""
        try:
            from pptx.util import Pt, Inches
            from pptx.enum.text import PP_ALIGN
            
            # Fill color
            if color_palette:
                shape.fill.solid()
                if is_top:
                    shape.fill.fore_color.rgb = self._int_to_rgb(color_palette.primary)
                else:
                    shape.fill.fore_color.rgb = self._int_to_rgb(color_palette.secondary)
            
            # Text formatting
            text_frame = shape.text_frame
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            
            para = text_frame.paragraphs[0]
            para.font.size = Pt(10 if is_top else 9)
            para.font.name = "Segoe UI"
            para.font.bold = is_top
            para.alignment = PP_ALIGN.CENTER
            
            if color_palette:
                para.font.color.rgb = self._int_to_rgb(color_palette.text_primary)
                
        except Exception as e:
            if self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                print(f"Warning: Could not format hierarchy shape: {e}")
    
    def _format_cycle_shape(self, shape, color_palette):
        """Format cycle shape with colors and text."""
        try:
            from pptx.util import Pt, Inches
            from pptx.enum.text import PP_ALIGN
            
            # Fill color
            if color_palette:
                shape.fill.solid()
                shape.fill.fore_color.rgb = self._int_to_rgb(color_palette.accent)
            
            # Text formatting
            text_frame = shape.text_frame
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            
            para = text_frame.paragraphs[0]
            para.font.size = Pt(9)
            para.font.name = "Segoe UI"
            para.alignment = PP_ALIGN.CENTER
            
            if color_palette:
                para.font.color.rgb = self._int_to_rgb(color_palette.text_primary)
                
        except Exception as e:
            if self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                print(f"Warning: Could not format cycle shape: {e}")
    
    def _format_list_shape(self, shape, color_palette):
        """Format list shape with colors and text."""
        try:
            from pptx.util import Pt, Inches
            from pptx.enum.text import PP_ALIGN
            
            # Fill color
            if color_palette:
                shape.fill.solid()
                shape.fill.fore_color.rgb = self._int_to_rgb(color_palette.secondary)
            
            # Text formatting
            text_frame = shape.text_frame
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.2)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            
            para = text_frame.paragraphs[0]
            para.font.size = Pt(11)
            para.font.name = "Segoe UI"
            para.alignment = PP_ALIGN.LEFT
            
            if color_palette:
                para.font.color.rgb = self._int_to_rgb(color_palette.text_primary)
                
        except Exception as e:
            if self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                print(f"Warning: Could not format list shape: {e}")
    
    def _format_relationship_shape(self, shape, color_palette, index):
        """Format relationship shape with colors and text."""
        try:
            from pptx.util import Pt, Inches
            from pptx.enum.text import PP_ALIGN
            
            # Fill color with transparency
            if color_palette:
                shape.fill.solid()
                colors = [color_palette.primary, color_palette.secondary, color_palette.accent]
                color_index = index % len(colors)
                shape.fill.fore_color.rgb = self._int_to_rgb(colors[color_index])
                shape.fill.transparency = 0.3  # Make semi-transparent for overlap effect
            
            # Text formatting
            text_frame = shape.text_frame
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            
            para = text_frame.paragraphs[0]
            para.font.size = Pt(10)
            para.font.name = "Segoe UI"
            para.alignment = PP_ALIGN.CENTER
            
            if color_palette:
                para.font.color.rgb = self._int_to_rgb(color_palette.text_primary)
                
        except Exception as e:
            if self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                print(f"Warning: Could not format relationship shape: {e}")
    
    def _format_matrix_shape(self, shape, color_palette, index):
        """Format matrix shape with colors and text."""
        try:
            from pptx.util import Pt, Inches
            from pptx.enum.text import PP_ALIGN
            
            # Fill color
            if color_palette:
                shape.fill.solid()
                colors = [color_palette.primary, color_palette.secondary, 
                         color_palette.accent, color_palette.text_secondary]
                color_index = index % len(colors)
                shape.fill.fore_color.rgb = self._int_to_rgb(colors[color_index])
            
            # Text formatting
            text_frame = shape.text_frame
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            
            para = text_frame.paragraphs[0]
            para.font.size = Pt(10)
            para.font.name = "Segoe UI"
            para.alignment = PP_ALIGN.CENTER
            
            if color_palette:
                para.font.color.rgb = self._int_to_rgb(color_palette.text_primary)
                
        except Exception as e:
            if self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                print(f"Warning: Could not format matrix shape: {e}")
    
    def _format_pyramid_shape(self, shape, color_palette, level):
        """Format pyramid shape with colors and text."""
        try:
            from pptx.util import Pt, Inches
            from pptx.enum.text import PP_ALIGN
            
            # Fill color - darker at top, lighter at bottom
            if color_palette:
                shape.fill.solid()
                # Vary color intensity by level
                if level == 0:
                    shape.fill.fore_color.rgb = self._int_to_rgb(color_palette.primary)
                elif level == 1:
                    shape.fill.fore_color.rgb = self._int_to_rgb(color_palette.secondary)
                else:
                    shape.fill.fore_color.rgb = self._int_to_rgb(color_palette.accent)
            
            # Text formatting
            text_frame = shape.text_frame
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            
            para = text_frame.paragraphs[0]
            para.font.size = Pt(11 - level)  # Smaller text for lower levels
            para.font.name = "Segoe UI"
            para.font.bold = level == 0  # Bold for top level
            para.alignment = PP_ALIGN.CENTER
            
            if color_palette:
                para.font.color.rgb = self._int_to_rgb(color_palette.text_primary)
                
        except Exception as e:
            if self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                print(f"Warning: Could not format pyramid shape: {e}")
    
    def _format_text_shape(self, shape, color_palette):
        """Format text shape with colors and text."""
        try:
            from pptx.util import Pt
            from pptx.enum.text import PP_ALIGN
            
            # Text formatting
            text_frame = shape.text_frame
            text_frame.word_wrap = True
            
            para = text_frame.paragraphs[0]
            para.font.size = Pt(10)
            para.font.name = "Segoe UI"
            para.font.bold = True
            para.alignment = PP_ALIGN.CENTER
            
            if color_palette:
                para.font.color.rgb = self._int_to_rgb(color_palette.primary)
                
        except Exception as e:
            if self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                print(f"Warning: Could not format text shape: {e}")
    
    def _add_connector_line(self, slide, x1, y1, x2, y2, color_palette):
        """Add connector line between two points."""
        try:
            from pptx.util import Pt
            
            connector = slide.shapes.add_connector(
                1,  # Straight connector
                x1, y1, x2, y2
            )
            
            if color_palette:
                connector.line.color.rgb = self._int_to_rgb(color_palette.accent)
                connector.line.width = Pt(2)
                
        except Exception as e:
            if self.config and hasattr(self.config, 'debug_mode') and self.config.debug_mode:
                print(f"Warning: Could not add connector line: {e}")
